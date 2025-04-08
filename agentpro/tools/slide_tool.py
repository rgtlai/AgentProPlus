from pptx import Presentation
from typing import List, Dict
import json
from .base import Tool
class SlideGenerationTool(Tool):
    name: str = "Slide Generation Tool"
    description: str = "A tool that can create a PPTX deck for a content. It takes a list of dictionaries. Each list dictionary item represents a slide in the presentation. Each dictionary item must have two keys: 'slide_title' and 'content'."
    arg: str = "List[Dict[slide_title, content]]. Ensure the Action Input is JSON parseable so I can convert it to required format"
    def run(self, slide_content: List[Dict[str, str]]) -> str:
        print(f"Calling Slide Generation Tool with slide_content TYPE :{type(slide_content)}")
        if type(slide_content) == str:
            try:
                slide_content = json.loads(slide_content)
                print("Converted Slide Content from str to JSON Dictionary")
            except Exception as e:
                return f"Error: {e}"    
        presentation = Presentation()
        # OPTIONAL : VARIABLE FONTS
        # OPTIONAL : TEXT COLORS
        # OPTIONAL : IMAGES / TABLES
        # Iterate over the slides list and add content to the presentation
        for slide in slide_content:    
            # Add a slide with a title and content layout
            slide_layout = presentation.slide_layouts[1]  # Layout 1 is 'Title and Content'
            ppt_slide = presentation.slides.add_slide(slide_layout)
            # Set the title and content for the slide
            title = ppt_slide.shapes.title
            content = ppt_slide.placeholders[1]
            title.text = slide['slide_title']
            content.text = slide['content']
        # Save the presentation to the specified output file
        output_path = "presentation.pptx"
        presentation.save(output_path)
        return f"Presentation saved as '{output_path}'."
